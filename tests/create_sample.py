"""Create a sample PPTX and dummy SRT for testing."""
import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import struct, wave

# ── Create sample PPTX ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Emu(9144000)   # 16:9 = 10 inches
prs.slide_height = Emu(5143500)   # 5.625 inches

slide_layout = prs.slide_layouts[1]  # title + content

def add_notes(slide, text: str):
    """슬라이드 노트 추가 헬퍼."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

# Slide 1
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "PPT2SlideDeck 데모"
title.text_frame.paragraphs[0].font.size = Pt(40)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1A, 0x53, 0xFF)
subtitle.text = "PPTX → 애니메이션 → MP4 자동 변환"
subtitle.text_frame.paragraphs[0].font.size = Pt(22)
add_notes(slide1, "안녕하세요!\nPPT2SlideDeck 데모를 시작하겠습니다.")

# Slide 2
slide2 = prs.slides.add_slide(slide_layout)
slide2.shapes.title.text = "주요 기능"
tf = slide2.placeholders[1].text_frame
tf.text = "타이핑 효과"
p2 = tf.add_paragraph(); p2.text = "페이드인 / 줌인 애니메이션"
p3 = tf.add_paragraph(); p3.text = "MP3 오디오 합성"
p4 = tf.add_paragraph(); p4.text = "자막(SRT) 삽입"
for para in tf.paragraphs:
    para.font.size = Pt(20)
add_notes(slide2, "슬라이드 노트에 입력한 텍스트가\n자동으로 자막이 됩니다.")

# Slide 3
slide3 = prs.slides.add_slide(slide_layout)
slide3.shapes.title.text = "완성!"
tf3 = slide3.placeholders[1].text_frame
tf3.text = "Claude Desktop MCP로 자동 실행됩니다."
tf3.paragraphs[0].font.size = Pt(22)
add_notes(slide3, "SRT 파일 없이도 노트만으로 자막이 생성됩니다!")

out_pptx = os.path.join(os.path.dirname(__file__), "fixtures", "sample.pptx")
os.makedirs(os.path.dirname(out_pptx), exist_ok=True)
prs.save(out_pptx)
print(f"Saved: {out_pptx}")

# ── Create dummy WAV (renamed .mp3 for test) ───────────────────────────────
# Generate a 10-second silent WAV file (ffmpeg handles it fine)
wav_path = os.path.join(os.path.dirname(__file__), "fixtures", "sample.mp3")
sample_rate = 44100
duration_sec = 10
n_samples = sample_rate * duration_sec
with wave.open(wav_path, "w") as wf:
    wf.setnchannels(1)
    wf.setsampwidth(2)
    wf.setframerate(sample_rate)
    wf.writeframes(b"\x00\x00" * n_samples)
print(f"Saved: {wav_path}")

# ── Create sample SRT ───────────────────────────────────────────────────────
srt_content = """\
1
00:00:00,000 --> 00:00:03,000
PPT2SlideDeck 데모 시작!

2
00:00:03,000 --> 00:00:07,000
주요 기능을 소개합니다.

3
00:00:07,000 --> 00:00:10,000
Claude Desktop MCP로 자동 변환 완료!
"""
srt_path = os.path.join(os.path.dirname(__file__), "fixtures", "sample.srt")
with open(srt_path, "w", encoding="utf-8") as f:
    f.write(srt_content)
print(f"Saved: {srt_path}")
print("\nAll sample files created!")
