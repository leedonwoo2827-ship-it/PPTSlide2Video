"""PPTX parsing logic using python-pptx."""
from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from slidecast.parser.slide_model import (
    ParagraphData, PresentationData, RunData, ShapeData, SlideData, TableCellData,
)
from slidecast.utils.color_utils import rgb_to_hex, theme_color_to_hex
from slidecast.utils.unit_utils import emu_to_px, pt_to_px

# Namespaces for animation XML parsing
_NSMAP = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}

ANIMATION_DELAY_STEP = 0.35  # seconds between shape animations


def parse_pptx(pptx_path: str, temp_dir: str) -> PresentationData:
    prs = Presentation(pptx_path)
    images_dir = os.path.join(temp_dir, "images")
    os.makedirs(images_dir, exist_ok=True)

    theme_map = _build_theme_color_map(prs)
    presentation = PresentationData(source_path=pptx_path, temp_dir=temp_dir)

    for idx, slide in enumerate(prs.slides):
        slide_data = _parse_slide(prs, slide, idx, images_dir, theme_map)
        presentation.slides.append(slide_data)

    return presentation


_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _build_theme_color_map(prs) -> dict[str, str]:
    """
    PPTX 테마 XML에서 scheme 색상 이름 → hex 값 맵을 구성한다.
    예: {"lt1": "#FFFFFF", "dk1": "#000000", "accent1": "#4472C4", ...}
    """
    result: dict[str, str] = {}
    try:
        master = prs.slide_masters[0]
        # 테마는 slide master's part의 theme XML에 있음
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        theme_elem = theme_part._element
        clr_elem = theme_elem.find(
            f".//{{{_A_NS}}}clrScheme"
        )
        if clr_elem is None:
            return result
        for child in clr_elem:
            name = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            # 각 색상 자식: <a:srgbClr val="RRGGBB"/> or <a:sysClr lastClr="RRGGBB"/>
            srgb = child.find(f"{{{_A_NS}}}srgbClr")
            if srgb is not None:
                result[name] = f"#{srgb.get('val', '000000').upper()}"
                continue
            sysclr = child.find(f"{{{_A_NS}}}sysClr")
            if sysclr is not None:
                result[name] = f"#{sysclr.get('lastClr', '000000').upper()}"
    except Exception:
        pass
    # fallback: Office 기본 테마 색상
    defaults = {
        "lt1": "#FFFFFF", "dk1": "#000000",
        "lt2": "#E7E6E6", "dk2": "#44546A",
        "accent1": "#4472C4", "accent2": "#ED7D31",
        "accent3": "#A9D18E", "accent4": "#FF0000",
        "accent5": "#FFC000", "accent6": "#70AD47",
    }
    for k, v in defaults.items():
        result.setdefault(k, v)
    return result


def _resolve_color_from_xml(run_elem, theme_map: dict[str, str]) -> str | None:
    """
    run XML element (<a:r>) 에서 텍스트 색상 추출.
    <a:rPr> → <a:solidFill> → srgbClr 또는 schemeClr → 테마 맵 조회.
    """
    if run_elem is None:
        return None
    # <a:rPr> (run properties) 탐색
    rPr = run_elem.find(f"{{{_A_NS}}}rPr")
    if rPr is None:
        return None
    solidFill = rPr.find(f"{{{_A_NS}}}solidFill")
    if solidFill is None:
        return None
    srgb = solidFill.find(f"{{{_A_NS}}}srgbClr")
    if srgb is not None:
        return f"#{srgb.get('val', '000000').upper()}"
    scheme = solidFill.find(f"{{{_A_NS}}}schemeClr")
    if scheme is not None:
        key = scheme.get("val", "")
        return theme_map.get(key)
    return None


def _parse_slide(prs, slide, idx: int, images_dir: str, theme_map: dict[str, str]) -> SlideData:
    width_px = emu_to_px(prs.slide_width)
    height_px = emu_to_px(prs.slide_height)

    bg_color = _extract_background(slide)
    animation_hints = _extract_animation_hints(slide)

    shapes: list[ShapeData] = []
    for z_order, shape in enumerate(slide.shapes):
        shape_data = _parse_shape(shape, z_order, images_dir, animation_hints, theme_map)
        if shape_data:
            shapes.append(shape_data)

    # Assign delays based on z_order (visible shapes sorted by order)
    visible = [s for s in shapes if s.animation_hint != "none"]
    for i, s in enumerate(visible):
        s.delay = i * ANIMATION_DELAY_STEP

    # Slide duration = last animation delay + its duration + 0.5s hold
    max_anim_duration = 1.2  # max single animation
    duration = (
        (len(visible) - 1) * ANIMATION_DELAY_STEP + max_anim_duration + 0.5
        if visible else 2.0
    )

    notes = None
    try:
        notes_tf = slide.notes_slide.notes_text_frame
        notes = notes_tf.text.strip() or None
    except Exception:
        pass

    return SlideData(
        slide_index=idx,
        width_px=width_px,
        height_px=height_px,
        background_color=bg_color,
        shapes=shapes,
        notes=notes,
        duration_seconds=round(duration, 2),
    )


def _extract_background(slide) -> str:
    """슬라이드 배경색을 CSS hex 값으로 반환. 테마 색상도 XML에서 직접 추출."""
    # 방법 1: python-pptx API (직접 RGB 지정인 경우)
    try:
        fill = slide.background.fill
        if fill.type is not None:
            color = theme_color_to_hex(fill.fore_color)
            if color:
                return color
    except Exception:
        pass

    # 방법 2: XML에서 배경 색상 직접 파싱 (테마 색상 포함)
    try:
        bg_elem = slide._element.find(
            ".//{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
        )
        if bg_elem is not None:
            # <a:solidFill><a:srgbClr val="RRGGBB"/>
            srgb = bg_elem.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
            )
            if srgb is not None:
                val = srgb.get("val", "")
                if val:
                    return f"#{val.upper()}"

            # <a:solidFill><a:sysClr lastClr="RRGGBB"/>
            sysclr = bg_elem.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}sysClr"
            )
            if sysclr is not None:
                val = sysclr.get("lastClr", "")
                if val:
                    return f"#{val.upper()}"
    except Exception:
        pass

    # 방법 3: 슬라이드 레이아웃/마스터에서 배경색 상속
    try:
        for parent in (slide.slide_layout, slide.slide_layout.slide_master):
            bg_elem = parent._element.find(
                ".//{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
            )
            if bg_elem is not None:
                srgb = bg_elem.find(
                    ".//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
                )
                if srgb is not None:
                    val = srgb.get("val", "")
                    if val:
                        return f"#{val.upper()}"
    except Exception:
        pass

    return "#FFFFFF"


def _extract_animation_hints(slide) -> dict[int, str]:
    """Parse slide XML timing tree to extract animation hints per shape id."""
    hints: dict[int, str] = {}
    try:
        timing = slide._element.find(".//p:timing", _NSMAP)
        if timing is None:
            return hints
        for anim in timing.iter():
            tag = anim.tag.split("}")[-1] if "}" in anim.tag else anim.tag
            if tag == "animEffect":
                filter_val = anim.get("filter", "")
                # Find parent cTn to get target shape id
                sp_tgt = anim.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}spTgt")
                if sp_tgt is not None:
                    sp_id = int(sp_tgt.get("spid", -1))
                    hints[sp_id] = _filter_to_hint(filter_val)
    except Exception:
        pass
    return hints


def _filter_to_hint(filter_val: str) -> str:
    f = filter_val.lower()
    if "wipe" in f or "fly" in f:
        return "fly_in_left"
    if "zoom" in f or "grow" in f:
        return "zoom_in"
    if "fade" in f or "dissolve" in f:
        return "fade_in"
    return "fade_in"


def _parse_shape(shape, z_order: int, images_dir: str,
                 animation_hints: dict[int, str],
                 theme_map: dict[str, str] | None = None) -> ShapeData | None:
    try:
        left = emu_to_px(shape.left or 0)
        top = emu_to_px(shape.top or 0)
        width = emu_to_px(shape.width or 0)
        height = emu_to_px(shape.height or 0)
        rotation = getattr(shape, "rotation", 0.0) or 0.0

        shape_type_str, anim_hint_default = _classify_shape(shape)
        anim_hint = animation_hints.get(shape.shape_id, anim_hint_default)

        fill_color = _extract_fill_color(shape)
        border_color = _extract_border_color(shape)
        border_width = _extract_border_width(shape)

        paragraphs: list[ParagraphData] = []
        text = None
        if shape.has_text_frame:
            text, paragraphs = _parse_text_frame(shape.text_frame, theme_map or {})

        image_path = None
        if shape_type_str == "picture":
            image_path = _extract_image(shape, images_dir)

        # 테이블 파싱
        table_data = None
        col_widths = None
        if shape_type_str == "table":
            table_data, col_widths = _parse_table(shape, theme_map or {})

        return ShapeData(
            shape_id=shape.shape_id,
            shape_type=shape_type_str,
            name=shape.name,
            left_px=left,
            top_px=top,
            width_px=width,
            height_px=height,
            rotation=rotation,
            text=text,
            paragraphs=paragraphs,
            fill_color=fill_color,
            border_color=border_color,
            border_width_px=border_width,
            image_path=image_path,
            table_data=table_data,
            col_widths=col_widths,
            z_order=z_order,
            animation_hint=anim_hint,
        )
    except Exception as e:
        return None


def _classify_shape(shape) -> tuple[str, str]:
    """Returns (shape_type_str, default_animation_hint)."""
    try:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            return "picture", "zoom_in"
        if st == MSO_SHAPE_TYPE.TABLE:
            return "table", "fade_in"
        if st == MSO_SHAPE_TYPE.GROUP:
            return "auto_shape", "fade_in"
        if st == MSO_SHAPE_TYPE.TEXT_BOX:
            return "text_box", "type_in"
    except Exception:
        pass
    # python-pptx: shape.has_table 체크
    try:
        if shape.has_table:
            return "table", "fade_in"
    except Exception:
        pass

    # Check by placeholder type
    try:
        ph = shape.placeholder_format
        if ph is not None:
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Pt
            ph_idx = ph.idx
            if ph_idx == 0:
                return "title", "type_in"
            return "text_box", "type_in"
    except Exception:
        pass

    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        return "text_box", "fade_in"

    return "auto_shape", "fade_in"


def _parse_text_frame(tf, theme_map: dict[str, str]) -> tuple[str, list[ParagraphData]]:
    paragraphs = []
    full_text_parts = []

    for para in tf.paragraphs:
        runs = []
        for run in para.runs:
            font = run.font
            color = "#000000"
            try:
                # Try XML-based resolution first (handles theme colors like lt1=white)
                c = _resolve_color_from_xml(run._r, theme_map)
                if not c:
                    c = theme_color_to_hex(font.color)
                if c:
                    color = c
            except Exception:
                pass

            font_size_px = 16.0
            try:
                if font.size:
                    font_size_px = pt_to_px(font.size.pt)
            except Exception:
                pass

            runs.append(RunData(
                text=run.text,
                bold=bool(font.bold),
                italic=bool(font.italic),
                font_size_px=font_size_px,
                color=color,
            ))

        align = "left"
        try:
            from pptx.enum.text import PP_ALIGN
            a = para.alignment
            if a == PP_ALIGN.CENTER:
                align = "center"
            elif a == PP_ALIGN.RIGHT:
                align = "right"
        except Exception:
            pass

        paragraphs.append(ParagraphData(runs=runs, alignment=align))
        full_text_parts.append("".join(r.text for r in runs))

    return "\n".join(full_text_parts), paragraphs


def _extract_fill_color(shape) -> str | None:
    # 1. python-pptx API 시도
    try:
        fill = shape.fill
        if fill.type is not None:
            c = theme_color_to_hex(fill.fore_color)
            if c:
                return c
    except Exception:
        pass
    # 2. XML에서 직접 추출 (spPr/solidFill/srgbClr)
    try:
        solid = shape._element.find(
            f".//{{{_A_NS}}}solidFill/{{{_A_NS}}}srgbClr"
        )
        if solid is not None:
            return f"#{solid.get('val', 'FFFFFF')}"
    except Exception:
        pass
    return None


def _extract_border_color(shape) -> str | None:
    try:
        line = shape.line
        if line.color and line.color.type is not None:
            return theme_color_to_hex(line.color)
    except Exception:
        pass
    return None


def _extract_border_width(shape) -> float:
    try:
        line = shape.line
        if line.width:
            return emu_to_px(line.width)
    except Exception:
        pass
    return 0.0


def _parse_table(shape, theme_map: dict[str, str]) -> tuple[list[list[TableCellData]], list[float]]:
    """테이블 shape에서 셀 데이터와 열 너비를 추출."""
    tbl = shape.table
    rows_data: list[list[TableCellData]] = []

    # 열 너비 (EMU → px)
    col_widths = [emu_to_px(col.width) for col in tbl.columns]

    for row in tbl.rows:
        row_cells: list[TableCellData] = []
        for cell in row.cells:
            # 셀 텍스트 파싱
            paragraphs: list[ParagraphData] = []
            text_parts: list[str] = []
            if cell.text_frame:
                text, paragraphs = _parse_text_frame(cell.text_frame, theme_map)
                text_parts.append(text or "")

            # 셀 배경색
            fill_color = None
            try:
                tc_xml = cell._tc
                solid = tc_xml.find(f".//{{{_A_NS}}}solidFill")
                if solid is not None:
                    srgb = solid.find(f"{{{_A_NS}}}srgbClr")
                    if srgb is not None:
                        fill_color = f"#{srgb.get('val', 'FFFFFF')}"
                    scheme = solid.find(f"{{{_A_NS}}}schemeClr")
                    if scheme is not None and not fill_color:
                        fill_color = theme_map.get(scheme.get("val", ""))
            except Exception:
                pass

            row_cells.append(TableCellData(
                text="".join(text_parts),
                paragraphs=paragraphs,
                fill_color=fill_color,
            ))
        rows_data.append(row_cells)

    return rows_data, col_widths


def _extract_image(shape, images_dir: str) -> str | None:
    try:
        img = shape.image
        ext = img.ext or "png"
        fname = f"img_{shape.shape_id}.{ext}"
        fpath = os.path.join(images_dir, fname)
        with open(fpath, "wb") as f:
            f.write(img.blob)
        return fpath
    except Exception:
        return None
